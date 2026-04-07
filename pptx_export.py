"""
38DN Pricing Model Review — PowerPoint Export
Generates branded PPTX output with comparison tables and summary slides.
Uses the 38DN "38N" theme: Century Gothic, navy/green/teal palette.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import io

# ---------------------------------------------------------------------------
# Brand constants (from 38N theme)
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x05, 0x0D, 0x25)
NAVY2     = RGBColor(0x21, 0x2B, 0x48)
GREEN     = RGBColor(0x45, 0xA7, 0x50)
TEAL      = RGBColor(0x51, 0x84, 0x84)
BLUE      = RGBColor(0x1D, 0x6F, 0xA9)
CYAN      = RGBColor(0x36, 0xAF, 0xCE)
LIGHT_GREY = RGBColor(0xE2, 0xE7, 0xF1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xB8, 0x32, 0x30)
TEXT_DARK = RGBColor(0x28, 0x28, 0x28)
GREEN_BG  = RGBColor(0xE2, 0xEF, 0xDA)
BLUE_BG   = RGBColor(0xE2, 0xE7, 0xF1)
WARN_BG   = RGBColor(0xE2, 0xE7, 0xF1)
RED_BG    = RGBColor(0xFD, 0xE2, 0xE1)

FONT_NAME = "Century Gothic"

# Slide dimensions (38DN custom widescreen)
SLIDE_W = 13817600
SLIDE_H = 7772400

# Layout positions
TITLE_LEFT   = Emu(811108)
TITLE_TOP    = Emu(335521)
TITLE_W      = Emu(11682483)
TITLE_H      = Emu(521127)
TABLE_LEFT   = Emu(811108)
TABLE_TOP    = Emu(1200000)
TABLE_W      = Emu(12200000)
FOOTER_LEFT  = Emu(535445)
FOOTER_TOP   = Emu(7559747)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def safe_float(v):
    if v is None: return None
    try: return float(v)
    except (ValueError, TypeError): return None


def fmt_neg(v, is_pct=False):
    """Format with parentheses for negatives."""
    if v is None: return "\u2014"
    if is_pct:
        return f"({abs(v):.2%})" if v < 0 else f"{v:.2%}"
    if abs(v) > 1000:
        return f"({abs(v):,.0f})" if v < 0 else f"{v:,.0f}"
    return f"({abs(v):.4f})" if v < 0 else f"{v:.4f}"


def set_cell_text(cell, text, font_size=10, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT):
    """Set cell text with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Cell margins
    cell.margin_left = Emu(71438)
    cell.margin_right = Emu(36000)
    cell.margin_top = Emu(18000)
    cell.margin_bottom = Emu(18000)


def set_cell_fill(cell, color):
    """Set solid fill on a cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_decorative_bars(slide):
    """Add the navy/green accent bars matching the 38DN template."""
    # Left nav bar
    left_bar = slide.shapes.add_shape(
        1, Emu(1), Emu(2671765), Emu(693018), Emu(2428875))  # 1 = rectangle
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = NAVY2
    left_bar.line.fill.background()

    # Right accent bar
    right_bar = slide.shapes.add_shape(
        1, Emu(13610398), Emu(2671764), Emu(232104), Emu(2428875))
    right_bar.fill.solid()
    right_bar.fill.fore_color.rgb = NAVY2
    right_bar.line.fill.background()


def add_title(slide, text):
    """Add a title text box in the 38DN style."""
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_NAME
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY2


def add_subtitle(slide, text, top=None):
    """Add a smaller subtitle."""
    t = top or Emu(900000)
    txBox = slide.shapes.add_textbox(TITLE_LEFT, t, TITLE_W, Emu(300000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(14)
    run.font.color.rgb = TEAL


def add_footer(slide, text="CONFIDENTIAL"):
    """Add footer text."""
    txBox = slide.shapes.add_textbox(FOOTER_LEFT, FOOTER_TOP, Emu(2000000), Emu(159488))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(8)
    run.font.color.rgb = TEAL


def add_table(slide, headers, rows, top=None, col_widths=None):
    """Add a formatted table to a slide.

    Args:
        headers: list of header strings
        rows: list of dicts, each with keys matching headers + optional '_styles' dict
              _styles = {col_idx: {"color": RGBColor, "fill": RGBColor, "bold": bool}}
        col_widths: optional list of EMU widths per column
    """
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    table_top = top or TABLE_TOP

    # Calculate available height
    max_height = Emu(6200000)
    row_height = min(Emu(280000), max_height // n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, TABLE_LEFT, table_top, TABLE_W, row_height * n_rows)
    table = table_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, NAVY)
        set_cell_text(cell, hdr, font_size=10, bold=True, color=WHITE)

    # Data rows
    for ri, row_data in enumerate(rows):
        styles = row_data.get("_styles", {})
        is_odd = ri % 2 == 1

        for ci, hdr in enumerate(headers):
            cell = table.cell(ri + 1, ci)
            val = row_data.get(hdr, "")
            style = styles.get(ci, {})

            # Background
            fill_color = style.get("fill")
            if fill_color:
                set_cell_fill(cell, fill_color)
            elif is_odd:
                set_cell_fill(cell, LIGHT_GREY)

            # Text
            text_color = style.get("color", TEXT_DARK)
            bold = style.get("bold", False)
            align = style.get("align", PP_ALIGN.LEFT)

            # Auto-detect negative formatting
            if isinstance(val, str) and val.startswith("(") and val.endswith(")"):
                text_color = RED
                bold = True

            set_cell_text(cell, val, font_size=10, bold=bold, color=text_color, alignment=align)

    return table


# ---------------------------------------------------------------------------
# Row categorization for slide breakout
# ---------------------------------------------------------------------------

# Which model rows belong to which category
CATEGORY_ROWS = {
    "System Details": [4, 7, 8, 10, 11, 12, 14, 15, 16, 18, 19, 21, 23, 24, 25],
    "CapEx": [117, 118, 119, 121, 122, 124, 126],
    "OpEx": [225, 226, 228, 239, 245, 281, 296, 301],
    "Timing": [67, 68, 69, 70, 71, 72, 73, 74, 75, 76, 77, 78, 79, 80, 143, 217],
    "Rates": [147, 148, 149, 150, 151, 152, 155, 156, 157, 158, 160, 161, 162,
              165, 166, 167, 168, 170],
    "Incentives & Tax": [216, 218, 591, 596, 597, 602],
}

# Flatten to see which rows are categorized
_CATEGORIZED = set()
for rows in CATEGORY_ROWS.values():
    _CATEGORIZED.update(rows)


def categorize_row(row_num):
    """Return the category for a given row number."""
    for cat, rows in CATEGORY_ROWS.items():
        if row_num in rows:
            return cat
    return "Other"


# ---------------------------------------------------------------------------
# Heat map color mapping
# ---------------------------------------------------------------------------

def heat_color(pct_diff):
    """Return a background color based on % difference magnitude."""
    if pct_diff is None:
        return None
    absp = abs(pct_diff)
    if absp > 0.10:     return RED_BG     # >10% difference
    elif absp > 0.05:   return WARN_BG    # 5-10%
    elif absp > 0.02:   return BLUE_BG    # 2-5%
    elif absp > 0:      return LIGHT_GREY  # <2%
    return None


# ---------------------------------------------------------------------------
# Main export function
# ---------------------------------------------------------------------------

TEXT_ROWS = {4, 8, 10, 18, 19, 21, 22, 117, 155, 156,
             165, 166, 591, 596}
DATE_ROWS = {68, 69, 70, 71, 72, 73}
PCT_ROWS = {30, 31, 36, 37, 158, 161, 162, 168, 219, 220, 221, 227, 229, 231,
            237, 241, 282, 283, 286, 288, 293, 297, 299, 597, 602}
DPW_ROWS = {32, 33, 38, 118, 119, 120, 121, 122, 123, 124, 126, 129, 157, 167, 216, 218}
INT_ROWS = {15, 16, 17, 24, 25, 143, 160, 170, 217, 225, 226, 228, 230, 235,
            240, 256, 258, 284, 285, 292, 296, 298, 302}

INPUT_ROW_LABELS = {
    4: "Project Name", 7: "Toggle (On/Off)", 8: "Developer Toggle",
    10: "Developer", 11: "Size MWDC", 12: "Size MWAC", 13: "DC-AC Ratio",
    14: "Energy Yield (kWh/WDC)", 15: "Year COD", 16: "System Life",
    17: "Effective System Life", 18: "State", 19: "Utility",
    21: "System Type", 22: "Customer",
    23: "Battery Storage", 24: "BESS Capacity kWac", 25: "BESS Capacity kWh",
    30: "FMV WACC (Target)", 31: "Live Appraisal IRR", 36: "Target IRR", 37: "Live Levered Pre-Tax IRR",
    68: "MIPA Signing", 69: "Close", 70: "NTP", 71: "MC", 72: "COD/PIS", 73: "SC",
    76: "MIPA->Close (mo)", 77: "Close->NTP (mo)", 78: "NTP->MC (mo)", 80: "COD->SC (mo)",
    117: "EPC Spend Curve", 118: "PV EPC Cost", 119: "PV LNTP Cost",
    120: "PV EPC + LNTP Cost", 121: "Customer Acquisition", 122: "IX Cost",
    123: "Closing and Legal", 124: "Other Capex Costs", 126: "ESS EPC Cost",
    129: "Total Capex Excl. Financing",
    143: "Revenue Lag", 147: "Rate Comp 1", 148: "Rate Comp 2",
    155: "Rate Name", 156: "Custom/Generic",
    157: "Energy Rate (at COD)", 158: "Energy Rate Escalator",
    160: "Rate Term", 161: "Rate Discount", 162: "Rate UCB Fee",
    165: "Rate 2 Name", 166: "Rate 2 Custom/Generic",
    167: "Rate 2 Energy Rate", 168: "Rate 2 Escalator", 170: "Rate 2 Term",
    216: "Upfront Incentive", 217: "Upfront Incentive Lag", 218: "ICSA Incentive",
    219: "ICSA COD %", 220: "ICSA Yr 1 %", 221: "ICSA Yr 2 %",
    225: "PV O&M Preventative", 226: "PV O&M Corrective", 227: "PV O&M Esc",
    228: "ESS O&M", 230: "AM Fee", 231: "AM Esc",
    234: "Upfront Bonus Lease", 235: "Lease Term", 236: "Lease (Year 1)", 237: "Lease Escalator",
    240: "Customer Mgmt Cost", 241: "Customer Mgmt Esc",
    255: "Inverter Replacement Toggle", 256: "Inverter Replacement $/MWac", 258: "Inverter Replacement Year",
    282: "Decom Disposal Cost %", 283: "Decom Cost Inflation", 284: "Decom End of Life",
    285: "Decom Bond $", 286: "Decom Annual Premium",
    291: "Custom PropTax Toggle", 292: "Property Taxes Yr 1", 293: "PropTax Escalator",
    296: "P&C Insurance", 297: "P&C Insurance Esc",
    298: "Catastrophic Coverage", 299: "Catastrophic Esc", 302: "Internal AM Costs",
    587: "COD Quarter",
    591: "Tax Treatment", 596: "TE Structure", 597: "ITC Rate", 602: "Eligible Costs %",
}

OUTPUT_ROWS = {32: "Dev Fee ($/W)", 33: "FMV Calculated ($/W)", 38: "NPP ($/W)", 39: "NPP ($)"}


def build_comparison_rows(proj1_data, proj2_data, m1_label, m2_label, rows_to_include=None):
    """Build comparison data for a set of rows."""
    all_rows = sorted(set(INPUT_ROW_LABELS.keys()) | set(OUTPUT_ROWS.keys()))
    if rows_to_include:
        all_rows = [r for r in all_rows if r in rows_to_include]

    result = []
    for r in all_rows:
        label = INPUT_ROW_LABELS.get(r, OUTPUT_ROWS.get(r, f"Row {r}"))
        v1_raw = proj1_data.get(r) if proj1_data else None
        v2_raw = proj2_data.get(r) if proj2_data else None

        if r in TEXT_ROWS or r in DATE_ROWS:
            result.append({
                "Row": str(r), "Field": label,
                m1_label: str(v1_raw) if v1_raw else "\u2014",
                m2_label: str(v2_raw) if v2_raw else "\u2014",
                "Delta": "", "\u0394 %": "",
                "_delta_raw": None, "_pct_raw": None,
            })
            continue

        v1 = safe_float(v1_raw)
        v2 = safe_float(v2_raw)
        is_pct = r in PCT_ROWS

        def fv(v):
            if v is None: return "\u2014"
            if is_pct: return f"{v:.2%}" if abs(v) < 1 else f"{v:.2f}%"
            if r in DPW_ROWS: return f"{v:.3f}"
            if r in INT_ROWS: return f"{v:,.0f}"
            if abs(v) > 1000: return f"{v:,.0f}"
            return f"{v:.4f}"

        delta = (v2 - v1) if v1 is not None and v2 is not None else None
        pct = (delta / abs(v1)) if delta is not None and v1 and v1 != 0 else None

        def fd(d, pct_fmt=False):
            if d is None: return "\u2014"
            if pct_fmt:
                return f"({abs(d):.2%})" if d < 0 else f"{d:.2%}"
            if d < 0:
                if is_pct: return f"({abs(d):.2%})"
                if abs(d) > 1000: return f"({abs(d):,.0f})"
                return f"({abs(d):.4f})"
            if is_pct: return f"{d:.2%}"
            if abs(d) > 1000: return f"{d:,.0f}"
            return f"{d:.4f}"

        result.append({
            "Row": str(r), "Field": label,
            m1_label: fv(v1), m2_label: fv(v2),
            "Delta": fd(delta), "\u0394 %": fd(pct, True),
            "_delta_raw": delta, "_pct_raw": pct,
        })

    return result


def generate_pptx(proj_name, proj1_data, proj2_data, m1_label, m2_label,
                  compare_bible=False, compare_model=False, bible_data=None,
                  include_summary_slide=True, include_category_slides=True):
    """Generate a branded 38DN PowerPoint with comparison analysis.

    Returns: BytesIO buffer containing the PPTX file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # ===================== TITLE SLIDE =====================
    slide = prs.slides.add_slide(blank_layout)

    # Full navy background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    # Green accent bar
    bar = slide.shapes.add_shape(1, Emu(0), Emu(5500000), Emu(SLIDE_W), Emu(400000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY2
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Emu(1500000), Emu(2200000), Emu(10000000), Emu(1200000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "38DN PRICING MODEL REVIEW"
    run.font.name = FONT_NAME
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = f"{proj_name}"
    run2.font.name = FONT_NAME
    run2.font.size = Pt(20)
    run2.font.color.rgb = WHITE

    # Comparison type
    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    parts = []
    if compare_bible: parts.append("vs. Pricing Bible Q1 2026")
    if compare_model: parts.append(f"{m1_label} vs. {m2_label}")
    run3.text = " | ".join(parts) if parts else "Model Review"
    run3.font.name = FONT_NAME
    run3.font.size = Pt(14)
    run3.font.color.rgb = TEAL

    add_footer(slide)

    # ===================== SUMMARY SLIDE =====================
    if include_summary_slide:
        slide = prs.slides.add_slide(blank_layout)
        add_decorative_bars(slide)
        add_title(slide, "SUMMARY \u2014 KEY DRIVERS")
        add_subtitle(slide, f"{proj_name} | {' | '.join(parts)}")

    # Build summary of key metrics
    key_outputs = [32, 33, 38, 39]  # Dev Fee, FMV, NPP
    key_inputs = [11, 12, 14, 118, 119, 122, 157, 158, 160, 216, 597, 602]
    summary_rows_list = build_comparison_rows(proj1_data, proj2_data, m1_label, m2_label,
                                              rows_to_include=key_outputs + key_inputs)

    if include_summary_slide:
        # Filter to rows that actually have deltas
        summary_with_delta = [r for r in summary_rows_list if r.get("_delta_raw") is not None and r["_delta_raw"] != 0]
        headers = ["Field", m1_label, m2_label, "Delta", "\u0394 %"]

        table_rows = []
        # Show outputs first (Dev Fee, FMV, NPP)
        for row in summary_rows_list:
            if int(row["Row"]) in key_outputs:
                pct = row.get("_pct_raw")
                styles = {}
                hc = heat_color(pct)
                if hc:
                    styles[3] = {"fill": hc}
                    styles[4] = {"fill": hc}
                if row.get("_delta_raw") is not None and row["_delta_raw"] < 0:
                    styles[3] = {**styles.get(3, {}), "color": RED}
                    styles[4] = {**styles.get(4, {}), "color": RED}
                table_rows.append({
                    "Field": row["Field"], m1_label: row[m1_label], m2_label: row[m2_label],
                    "Delta": row["Delta"], "\u0394 %": row["\u0394 %"],
                    "_styles": {0: {"bold": True, "fill": NAVY2, "color": WHITE}, **styles},
                })

        # Then key inputs with changes
        for row in summary_with_delta:
            if int(row["Row"]) not in key_outputs:
                pct = row.get("_pct_raw")
                styles = {}
                hc = heat_color(pct)
                if hc:
                    styles[3] = {"fill": hc}
                    styles[4] = {"fill": hc}
                if row["_delta_raw"] < 0:
                    styles[3] = {**styles.get(3, {}), "color": RED}
                    styles[4] = {**styles.get(4, {}), "color": RED}
                table_rows.append({
                    "Field": row["Field"], m1_label: row[m1_label], m2_label: row[m2_label],
                    "Delta": row["Delta"], "\u0394 %": row["\u0394 %"],
                    "_styles": styles,
                })

        if table_rows:
            add_table(slide, headers, table_rows,
                      col_widths=[3500000, 2500000, 2500000, 2000000, 1700000])

        add_footer(slide)

    # ===================== CATEGORY SLIDES =====================
    if include_category_slides:
        categories_order = ["System Details", "CapEx", "OpEx", "Timing", "Rates", "Incentives & Tax", "Other"]

        all_comp_rows = build_comparison_rows(proj1_data, proj2_data, m1_label, m2_label)

        for cat in categories_order:
            cat_rows_nums = CATEGORY_ROWS.get(cat, None)

            if cat == "Other":
                # Everything not in other categories
                cat_data = [r for r in all_comp_rows if int(r["Row"]) not in _CATEGORIZED
                            and int(r["Row"]) not in [32, 33, 38, 39]]
            else:
                if not cat_rows_nums:
                    continue
                cat_data = [r for r in all_comp_rows if int(r["Row"]) in cat_rows_nums]

            # Only include rows that have data or differences
            cat_data_filtered = [r for r in cat_data
                                 if r[m1_label] != "\u2014" or r[m2_label] != "\u2014"]

            if not cat_data_filtered:
                continue

            slide = prs.slides.add_slide(blank_layout)
            add_decorative_bars(slide)
            add_title(slide, f"{cat.upper()} INPUTS")
            add_subtitle(slide, proj_name)

            table_rows = []
            for row in cat_data_filtered:
                pct = row.get("_pct_raw")
                styles = {}
                hc = heat_color(pct)
                if hc:
                    styles[4] = {"fill": hc}
                    styles[5] = {"fill": hc}
                if row.get("_delta_raw") is not None and row["_delta_raw"] < 0:
                    styles[4] = {**styles.get(4, {}), "color": RED}
                    styles[5] = {**styles.get(5, {}), "color": RED}

                table_rows.append({
                    "Row": row["Row"], "Field": row["Field"],
                    m1_label: row[m1_label], m2_label: row[m2_label],
                    "Delta": row["Delta"], "\u0394 %": row["\u0394 %"],
                    "_styles": styles,
                })

            # Paginate if too many rows (max ~18 per slide)
            max_per_slide = 18
            for page_start in range(0, len(table_rows), max_per_slide):
                page_rows = table_rows[page_start:page_start + max_per_slide]
                if page_start > 0:
                    slide = prs.slides.add_slide(blank_layout)
                    add_decorative_bars(slide)
                    add_title(slide, f"{cat.upper()} INPUTS (CONT.)")
                    add_subtitle(slide, proj_name)

                add_table(slide, ["Row", "Field", m1_label, m2_label, "Delta", "\u0394 %"],
                          page_rows,
                          col_widths=[700000, 3200000, 2300000, 2300000, 1900000, 1800000])
                add_footer(slide)

    # ===================== Save to buffer =====================
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
